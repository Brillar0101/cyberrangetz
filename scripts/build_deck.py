from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors (Airbnb-clean style, TZ accent) ────
GREEN      = RGBColor(0x1E, 0xB5, 0x3A)   # TZ green — primary accent
BLUE       = RGBColor(0x00, 0xA3, 0xDD)   # TZ blue
YELLOW     = RGBColor(0xFC, 0xD1, 0x16)   # TZ yellow
DARK       = RGBColor(0x48, 0x48, 0x48)   # body text
BLACK      = RGBColor(0x2C, 0x2C, 0x2C)   # headings
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x99, 0x99, 0x99)
LIGHT      = RGBColor(0xF5, 0xF5, 0xF5)   # subtle bg
BORDER     = RGBColor(0xDD, 0xDD, 0xDD)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def box(sl, l, t, w, h, fill=None, line=None):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def t(sl, text, l, t_, w, h, sz=12, bold=False, color=DARK,
      align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t_), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Helvetica"
    return tb

def img(sl, l, t_, w, h, label="[IMAGE]"):
    box(sl, l, t_, w, h, fill=WHITE, line=BORDER)
    t(sl, label, l, t_+h/2-0.2, w, 0.4, sz=9, color=BORDER, align=PP_ALIGN.CENTER, italic=True)

def heading(sl, title):
    t(sl, title, 0.8, 0.6, 8, 0.7, sz=36, bold=True, color=GREEN)
    box(sl, 0.8, 1.35, 11.7, 0.05, fill=GREEN)

def border_slide():
    sl = blank()
    box(sl, 0.15, 0.15, 13.03, 7.2, line=BORDER)
    return sl


# ══════════════════════════════════════════════
# 1 — TITLE
# ══════════════════════════════════════════════
sl = blank()
img(sl, 0, 0, 6.5, 7.5,
    "[AI IMAGE]\nAfrican student at laptop\nDark terminal, dramatic lighting")
box(sl, 6.5, 0, 6.83, 7.5, fill=WHITE)

box(sl, 7.5, 2.0, 0.55, 0.55, fill=GREEN)
t(sl, "CyberRange Tanzania", 8.2, 1.95, 4.5, 0.6,
  sz=22, bold=True, color=GREEN)

t(sl, "Pitch Deck", 7.5, 3.8, 5.0, 0.6,
  sz=20, bold=True, color=BLACK)
t(sl, "Building Africa\u2019s cybersecurity workforce,\none lab at a time.",
  7.5, 4.5, 4.5, 0.8, sz=14, color=DARK)


# ══════════════════════════════════════════════
# 2 — PROBLEM
# ══════════════════════════════════════════════
sl = border_slide()
heading(sl, "Problem")

img(sl, 6.8, 0.5, 6.2, 6.7,
    "[AI IMAGE]\nTanzanian student frustrated at laptop\nNo cybersecurity lab access")

t(sl, "67.1 million", 0.8, 2.0, 5.5, 0.5, sz=18, bold=True, color=GREEN)
t(sl, "telecom subscriptions in Tanzania, growing 7.8%/year.\nYet cybercrime surged 150% in a single year.",
  0.8, 2.55, 5.5, 0.8, sz=13, color=DARK)

t(sl, "102", 0.8, 3.7, 5.5, 0.5, sz=18, bold=True, color=GREEN)
t(sl, "ICT security personnel in the entire\nTanzanian public service.",
  0.8, 4.25, 5.5, 0.7, sz=13, color=DARK)

t(sl, "<300,000", 0.8, 5.2, 5.5, 0.5, sz=18, bold=True, color=GREEN)
t(sl, "cyber professionals protect all of Africa\u2019s\n1.4 billion people. Only 20,000 are certified.",
  0.8, 5.75, 5.5, 0.7, sz=13, color=DARK)


# ══════════════════════════════════════════════
# 3 — NON-OBVIOUS INSIGHT
# ══════════════════════════════════════════════
sl = border_slide()
heading(sl, "The Insight")

img(sl, 6.8, 0.5, 6.2, 6.7,
    "[AI IMAGE]\nTanzania government building\nor ITU policy visual")

t(sl, "Tanzania ranks #2 in Africa\non the ITU Global Cybersecurity Index.",
  0.8, 2.2, 5.5, 0.8, sz=18, bold=True, color=BLACK)

t(sl, "Score: 90.58 / 100", 0.8, 3.2, 5.5, 0.5, sz=14, color=GREEN)

t(sl, "Strong governance and policy.\nAlmost zero training infrastructure.",
  0.8, 4.0, 5.5, 0.8, sz=14, color=DARK)

t(sl, "The government\u2019s own Cybersecurity Strategy\n2022\u20132027 calls for \u201cone platform for cyber security\nawareness and knowledge testing\u201d \u2014\nbut has not built it.",
  0.8, 5.0, 5.5, 1.2, sz=13, color=GRAY)


# ══════════════════════════════════════════════
# 4 — SOLUTION
# ══════════════════════════════════════════════
sl = border_slide()
heading(sl, "Solution")

img(sl, 0.3, 4.5, 12.7, 2.8,
    "[AI IMAGE] Split-screen browser: lab instructions left, dark Linux terminal right")

cols = [
    ("Browser-Based Labs",
     "Real Linux terminal in the\nbrowser. Any device, any\nlocation, no setup."),
    ("Attack, Then Defend",
     "Every topic taught from both\nsides. Hack it first, then\nlearn to stop it."),
    ("Built for Africa",
     "M-Pesa fraud, SIM swapping,\ntelecom attacks. Real threats,\nnot Western scenarios."),
]
x = 0.8
for title, body in cols:
    img(sl, x + 0.7, 1.7, 2.2, 1.6, "[ICON]")
    t(sl, title, x, 3.45, 3.6, 0.4, sz=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    t(sl, body, x, 3.9, 3.6, 0.6, sz=11, color=DARK, align=PP_ALIGN.CENTER)
    x += 4.0


# ══════════════════════════════════════════════
# 5 — MARKET VALIDATION
# ══════════════════════════════════════════════
sl = border_slide()
heading(sl, "Market Validation")

# Two big circles
stats = [
    ("3.9M", "Global cybersecurity\nworkforce gap", "ISC\u00b2 Workforce Study 2024"),
    ("$4.12B", "Cybercrime cost to Africa\n(>10% of GDP)", "Cisco/Access Partnership"),
]
x = 1.5
for number, label, source in stats:
    box(sl, x, 2.2, 4.5, 3.5, fill=GREEN)
    t(sl, number, x, 2.8, 4.5, 1.0, sz=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    t(sl, label, x, 3.9, 4.5, 0.8, sz=14, color=WHITE, align=PP_ALIGN.CENTER)
    t(sl, source, x, 6.0, 4.5, 0.4, sz=9, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
    x += 5.5


# ══════════════════════════════════════════════
# 6 — MARKET SIZE
# ══════════════════════════════════════════════
sl = border_slide()
heading(sl, "Market Opportunity")

img(sl, 6.8, 0.5, 6.2, 6.7,
    "[AI IMAGE]\nMap of East Africa\nGlowing network nodes")

t(sl, "1,200\u20131,614", 0.8, 2.2, 5.5, 0.5, sz=24, bold=True, color=GREEN)
t(sl, "ICT students at UDSM CoICT. Our pilot.", 0.8, 2.8, 5.5, 0.4, sz=13, color=DARK)

t(sl, "3,050\u20134,714", 0.8, 3.6, 5.5, 0.5, sz=24, bold=True, color=GREEN)
t(sl, "ICT students across 6 Tanzanian universities.", 0.8, 4.2, 5.5, 0.4, sz=13, color=DARK)

t(sl, "$94,800/yr", 0.8, 5.0, 5.5, 0.5, sz=24, bold=True, color=GREEN)
t(sl, "Phase 2 revenue \u2014 6 university licenses.", 0.8, 5.6, 5.5, 0.4, sz=13, color=DARK)


# ══════════════════════════════════════════════
# 7 — PRODUCT (divider)
# ══════════════════════════════════════════════
sl = blank()
box(sl, 0, 0, 6.5, 7.5, fill=DARK)
t(sl, "Product", 0.8, 3.0, 5.5, 1.0, sz=48, bold=True, color=WHITE)
box(sl, 0.8, 4.1, 4.5, 0.06, fill=GREEN)

img(sl, 6.5, 0, 6.83, 7.5,
    "[AI IMAGE]\nMacBook showing CyberRange TZ\nSplit-screen terminal + instructions")


# ══════════════════════════════════════════════
# 8 — PRODUCT: Dashboard
# ══════════════════════════════════════════════
sl = border_slide()

img(sl, 3.5, 0.8, 9.3, 6.2,
    "[SCREENSHOT]\nCyberRange TZ dashboard\nModule cards with Red/Blue badges")

img(sl, 0.4, 4.2, 0.8, 0.8, "[ICON]")
box(sl, 0.3, 5.1, 2.8, 0.4, fill=GREEN)
t(sl, "Choose Your Module", 0.35, 5.1, 2.7, 0.4,
  sz=11, bold=True, color=WHITE)
t(sl, "5 modules. 54 labs.\nRed team first, then Blue.", 0.3, 5.6, 3.0, 0.7, sz=11, color=DARK)


# ══════════════════════════════════════════════
# 9 — PRODUCT: Terminal
# ══════════════════════════════════════════════
sl = border_slide()

img(sl, 3.5, 0.8, 9.3, 6.2,
    "[SCREENSHOT]\nSplit screen: instructions left\nLive Linux terminal with nmap output right")

img(sl, 0.4, 4.2, 0.8, 0.8, "[ICON]")
box(sl, 0.3, 5.1, 2.8, 0.4, fill=GREEN)
t(sl, "Real Linux Terminal", 0.35, 5.1, 2.7, 0.4,
  sz=11, bold=True, color=WHITE)
t(sl, "Docker container in the cloud.\nReal Kali shell in the browser.", 0.3, 5.6, 3.0, 0.7, sz=11, color=DARK)


# ══════════════════════════════════════════════
# 10 — PRODUCT: Red → Blue
# ══════════════════════════════════════════════
sl = border_slide()

img(sl, 3.5, 0.8, 9.3, 6.2,
    "[SCREENSHOT]\nRED badge 'Scan the Network' \u2014 completed\nBLUE badge 'Detect the Scanner' \u2014 unlocked")

img(sl, 0.4, 4.2, 0.8, 0.8, "[ICON]")
box(sl, 0.3, 5.1, 2.8, 0.4, fill=GREEN)
t(sl, "Attack, Then Defend", 0.35, 5.1, 2.7, 0.4,
  sz=11, bold=True, color=WHITE)
t(sl, "Complete the Red lab.\nBlue lab unlocks.", 0.3, 5.6, 3.0, 0.7, sz=11, color=DARK)


# ══════════════════════════════════════════════
# 11 — BUSINESS MODEL
# ══════════════════════════════════════════════
sl = border_slide()
heading(sl, "Business Model")

img(sl, 6.8, 0.5, 6.2, 6.7,
    "[AI IMAGE]\nTanzanian university campus\nStudents in modern computer lab")

t(sl, "Flat institutional license. Students access free.",
  0.8, 1.8, 5.5, 0.4, sz=12, color=GRAY, italic=True)

t(sl, "$2,500/mo", 0.8, 2.8, 5.5, 0.6, sz=32, bold=True, color=BLACK)
t(sl, "Premium tier \u2014 up to 1,500 students\nCost per student: $1.67/month",
  0.8, 3.5, 5.5, 0.6, sz=13, color=DARK)

t(sl, "$30,000", 0.8, 4.6, 5.5, 0.6, sz=32, bold=True, color=BLACK)
t(sl, "Year 1 revenue\nUDSM pilot only",
  0.8, 5.3, 5.5, 0.6, sz=13, color=DARK)


# ══════════════════════════════════════════════
# 12 — GO-TO-MARKET
# ══════════════════════════════════════════════
sl = border_slide()
heading(sl, "Go-to-Market")

channels = [
    ("[AI IMAGE]\nICT Commission\nGovernment meeting",
     "ICT Commission",
     "Direct pipeline to 50+\naccredited universities."),
    ("[AI IMAGE]\nTanzanian university\nstudents in CS lab",
     "University Pilots",
     "3 pilot universities.\nFree semester for feedback."),
    ("[AI IMAGE]\nGovernment employees\nin training session",
     "Gov Training",
     "500 personnel target.\nStrategy 2022 mandate."),
]
x = 0.5
for image, title, body in channels:
    img(sl, x, 1.7, 3.8, 2.8, image)
    t(sl, title, x, 4.7, 3.8, 0.4, sz=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    t(sl, body, x, 5.2, 3.8, 0.8, sz=12, color=DARK, align=PP_ALIGN.CENTER)
    x += 4.1


# ══════════════════════════════════════════════
# 13 — COMPETITION (divider)
# ══════════════════════════════════════════════
sl = blank()
box(sl, 0, 0, 6.5, 7.5, fill=WHITE)
t(sl, "Competition", 0.8, 3.0, 5.5, 1.0, sz=44, bold=True, color=GREEN)
box(sl, 0.8, 4.1, 5.0, 0.06, fill=GREEN)

img(sl, 6.5, 0, 6.83, 7.5,
    "[AI IMAGE]\nAbstract competitive landscape\nCyberRange TZ glowing center")


# ══════════════════════════════════════════════
# 14 — COMPETITIVE MATRIX
# ══════════════════════════════════════════════
sl = border_slide()

# Axis labels
t(sl, "Affordable", 5.7, 0.4, 2.0, 0.4, sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
t(sl, "Expensive", 5.7, 6.8, 2.0, 0.4, sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
t(sl, "Generic /\nGlobal", 0.1, 3.3, 1.5, 0.7, sz=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
t(sl, "Africa-\nFocused", 11.5, 3.3, 1.5, 0.7, sz=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)

# Crosshair
box(sl, 1.6, 3.65, 10.1, 0.04, fill=GREEN)
box(sl, 6.6, 0.9, 0.04, 5.8, fill=GREEN)

# Competitors
comps = [
    ("TryHackMe", 2.5, 2.0),
    ("Hack The Box", 2.5, 4.8),
    ("SANS", 3.5, 5.5),
    ("DarTU (offline)", 5.2, 4.5),
]
for name, cx, cy in comps:
    box(sl, cx, cy, 0.35, 0.35, fill=GRAY)
    t(sl, name, cx-0.5, cy+0.4, 2.0, 0.4, sz=10, color=GRAY, align=PP_ALIGN.CENTER)

# Us
box(sl, 9.5, 1.5, 0.4, 0.4, fill=GREEN)
t(sl, "CyberRange TZ", 8.5, 2.0, 2.5, 0.4, sz=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
t(sl, "You", 10.0, 1.5, 1.0, 0.4, sz=10, color=GREEN)


# ══════════════════════════════════════════════
# 15 — COMPETITIVE ADVANTAGE
# ══════════════════════════════════════════════
sl = border_slide()
heading(sl, "Competitive Advantage")

advantages = [
    ("First to Market",
     "No cloud cyber lab exists\nfor East African universities."),
    ("African Threats",
     "M-Pesa fraud, SIM swapping,\ntelecom attacks."),
    ("Purple Team",
     "Every attack lab paired\nwith a defense lab."),
    ("Gov Alignment",
     "Directly fulfills Strategy\n2022 training mandates."),
]
x = 0.5
for title, body in advantages:
    img(sl, x + 0.6, 2.0, 1.7, 1.7, "[ICON]")
    t(sl, title, x, 4.0, 2.9, 0.4, sz=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    t(sl, body, x, 4.5, 2.9, 0.8, sz=11, color=DARK, align=PP_ALIGN.CENTER)
    x += 3.1


# ══════════════════════════════════════════════
# 16 — TEAM
# ══════════════════════════════════════════════
sl = border_slide()
heading(sl, "Team")

team = [
    ("[YOUR PHOTO]", "Barakaeli Lawuo", "CTO & Co-Founder",
     "Computer Engineering, Virginia Tech.\nCybersecurity & Networks. Built SwishVision AI."),
    ("[AUDREY PHOTO]", "Audrey Lyimo", "COO & Co-Founder",
     "[One-line proof of relevant\noperations / business experience]"),
    ("[RICHARD PHOTO]", "Richard Kiwelu", "Chief Marketing Officer",
     "[One-line proof of relevant\nmarketing / growth experience]"),
]
x = 1.0
for photo, name, role, bio in team:
    img(sl, x + 0.3, 1.8, 2.8, 2.8, photo)
    t(sl, name, x, 4.8, 3.4, 0.4, sz=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    t(sl, role, x, 5.25, 3.4, 0.35, sz=11, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    t(sl, bio, x, 5.7, 3.4, 0.8, sz=10, color=GRAY, align=PP_ALIGN.CENTER)
    x += 3.8


# ══════════════════════════════════════════════
# 17 — TRACTION
# ══════════════════════════════════════════════
sl = border_slide()
heading(sl, "Traction")

img(sl, 6.8, 0.5, 6.2, 6.7,
    "[SCREENSHOT]\nWorking MVP: browser terminal\nRunning attack/defend lab")

t(sl, "MVP Built and Functional", 0.8, 2.0, 5.5, 0.4,
  sz=16, bold=True, color=BLACK)

t(sl, "Working browser-based Linux terminal\nwith attack & defend lab environment.\n\n54 labs across 5 modules.\nCustom honeypot targets deployed.\nMulti-container Docker networking.",
  0.8, 2.6, 5.5, 2.5, sz=13, color=DARK)

box(sl, 0.8, 5.5, 5.5, 0.05, fill=GREEN)
t(sl, "Next: UDSM CoICT pilot deployment\nand first institutional license.",
  0.8, 5.7, 5.5, 0.7, sz=12, color=GRAY, italic=True)


# ══════════════════════════════════════════════
# 18 — FINANCIAL
# ══════════════════════════════════════════════
sl = border_slide()
heading(sl, "Financial")

img(sl, 6.8, 1.5, 5.8, 5.0,
    "[CHART]\nDonut chart: Use of funds\n28% Platform MVP\n20% Honeypots + Content\n30% Co-founders\n12% Cloud Infra\n20% Marketing + Ops")

t(sl, "$50K Seed Round", 0.8, 2.0, 5.5, 0.6,
  sz=28, bold=True, color=BLACK)
t(sl, "12 months runway", 0.8, 2.7, 5.5, 0.4, sz=14, color=GRAY)

t(sl, "$30,000", 0.8, 3.6, 2.5, 0.5, sz=24, bold=True, color=GREEN)
t(sl, "Year 1 Revenue\nUDSM pilot", 3.5, 3.6, 3.0, 0.5, sz=12, color=DARK)

t(sl, "$94,800", 0.8, 4.4, 2.5, 0.5, sz=24, bold=True, color=GREEN)
t(sl, "Phase 2 Revenue\n6 universities", 3.5, 4.4, 3.0, 0.5, sz=12, color=DARK)

t(sl, "$28,400", 0.8, 5.2, 2.5, 0.5, sz=24, bold=True, color=GREEN)
t(sl, "Cash remaining\nend of Year 1", 3.5, 5.2, 3.0, 0.5, sz=12, color=DARK)

t(sl, "Near breakeven with UDSM alone.\nAdding a second university = profitable.",
  0.8, 6.2, 5.5, 0.6, sz=11, color=GRAY, italic=True)


# ══════════════════════════════════════════════
# 19 — CLOSING
# ══════════════════════════════════════════════
sl = blank()
box(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0x1A, 0x1A, 0x1A))

t(sl, "Raising $50,000 to launch\nbrowser-based attack & defend labs at UDSM\nand train 1,200+ cybersecurity students in Year 1.",
  1.0, 1.5, 11.0, 1.5, sz=20, color=WHITE, align=PP_ALIGN.CENTER)

box(sl, 5.0, 3.5, 3.3, 0.06, fill=GREEN)

t(sl, "Building Africa\u2019s cybersecurity workforce,\none lab at a time.",
  1.0, 4.0, 11.0, 1.0, sz=24, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

t(sl, "barakaelilawuo0101@gmail.com", 1.0, 5.8, 11.0, 0.5,
  sz=14, color=GRAY, align=PP_ALIGN.CENTER)

# TZ flag bar
box(sl, 0, 7.2, 3.33, 0.12, fill=GREEN)
box(sl, 3.33, 7.2, 3.33, 0.12, fill=YELLOW)
box(sl, 6.66, 7.2, 3.34, 0.12, fill=RGBColor(0x1A, 0x1A, 0x1A))
box(sl, 10.0, 7.2, 3.33, 0.12, fill=BLUE)


# ══════════════════════════════════════════════
prs.save("CyberRangeTZ_PitchDeck_v6.pptx")
print("Saved: CyberRangeTZ_PitchDeck_v6.pptx - 19 slides")
